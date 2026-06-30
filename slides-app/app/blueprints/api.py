from flask import Blueprint, request, jsonify, g, current_app, send_file
from flask_babel import _
import os
import uuid
from io import BytesIO
from werkzeug.utils import secure_filename
from app.repositories import presentation_repository, slide_repository, slide_component_repository
from app.auth import login_required

bp = Blueprint('api', __name__, url_prefix='/api')

SLIDE_TEMPLATES = {
    'blank': [],
    'title_text': [
        {'type': 'title', 'content': 'Titolo', 'x': 5, 'y': 5, 'width': 90, 'height': 20,
         'font_size': 48, 'color': '#ff6600', 'bg_color': 'transparent', 'z_index': 0, 'image': None},
        {'type': 'text', 'content': 'Testo della slide...', 'x': 5, 'y': 30, 'width': 90, 'height': 62,
         'font_size': 22, 'color': '#333333', 'bg_color': 'transparent', 'z_index': 0, 'image': None},
    ],
    'title_text_image': [
        {'type': 'title', 'content': 'Titolo', 'x': 5, 'y': 5, 'width': 90, 'height': 18,
         'font_size': 48, 'color': '#ff6600', 'bg_color': 'transparent', 'z_index': 0, 'image': None},
        {'type': 'text', 'content': 'Testo della slide...', 'x': 5, 'y': 28, 'width': 56, 'height': 65,
         'font_size': 20, 'color': '#333333', 'bg_color': 'transparent', 'z_index': 0, 'image': None},
        {'type': 'image', 'content': '', 'x': 63, 'y': 28, 'width': 32, 'height': 65,
         'font_size': 0, 'color': '#000000', 'bg_color': 'transparent', 'z_index': 0, 'image': None},
    ],
}


def save_image(file):
    if file and file.filename:
        filename = secure_filename(file.filename)
        upload_folder = current_app.config['UPLOAD_FOLDER']
        os.makedirs(upload_folder, exist_ok=True)
        filepath = os.path.join(upload_folder, filename)
        file.save(filepath)
        return f'/static/uploads/{filename}'
    return None


# ─── Presentations ────────────────────────────────────────────────────────────

@bp.route('/presentations')
@login_required
def get_presentations():
    presentations = presentation_repository.get_presentations_by_author(g.user['id'])
    return jsonify({'presentations': presentations})


@bp.route('/presentations/create', methods=['POST'])
@login_required
def create_presentation():
    data = request.get_json() or {}
    title = data.get('title', '').strip()
    description = data.get('description', '').strip()
    if not title:
        return jsonify({'ok': False, 'error': _('Il titolo è obbligatorio.')}), 400
    presentation = presentation_repository.create_presentation(title, description, g.user['id'])
    return jsonify({'ok': True, 'presentation': presentation})


@bp.route('/presentations/<int:presentation_id>', methods=['DELETE'])
@login_required
def delete_presentation(presentation_id):
    presentation = presentation_repository.get_presentation_by_id(presentation_id)
    if not presentation or presentation['author_id'] != g.user['id']:
        return jsonify({'ok': False, 'error': _('Presentazione non trovata.')}), 404
    presentation_repository.delete_presentation(presentation_id)
    return jsonify({'ok': True})


@bp.route('/presentations/<int:presentation_id>/slides')
@login_required
def get_slides(presentation_id):
    presentation = presentation_repository.get_presentation_by_id(presentation_id)
    if not presentation or presentation['author_id'] != g.user['id']:
        return jsonify({'ok': False, 'error': _('Presentazione non trovata.')}), 404
    slides = slide_repository.get_slides_by_presentation_id(presentation_id)
    return jsonify({'slides': slides})


@bp.route('/presentations/<int:presentation_id>/slides', methods=['POST'])
@login_required
def create_slide(presentation_id):
    presentation = presentation_repository.get_presentation_by_id(presentation_id)
    if not presentation or presentation['author_id'] != g.user['id']:
        return jsonify({'ok': False, 'error': _('Presentazione non trovata.')}), 404
    data = request.get_json() or {}
    bg_color = data.get('bg_color', '#ffffff')
    template_name = data.get('template', 'blank')
    slides = slide_repository.get_slides_by_presentation_id(presentation_id)
    position = len(slides) + 1
    slide = slide_repository.create_slide(presentation_id, position, bg_color)
    preset = SLIDE_TEMPLATES.get(template_name, [])
    if preset:
        slide_component_repository.save_all_components(slide['id'], preset)
    slide['components'] = slide_component_repository.get_components_by_slide(slide['id'])
    return jsonify({'ok': True, 'slide': slide})


# ─── Slides ───────────────────────────────────────────────────────────────────

@bp.route('/slides/<int:slide_id>/move_up', methods=['POST'])
@login_required
def move_slide_up(slide_id):
    slide = slide_repository.get_slide_by_id(slide_id)
    if not slide:
        return jsonify({'ok': False, 'error': 'Slide non trovata.'}), 404
    presentation = presentation_repository.get_presentation_by_id(slide['presentation_id'])
    if presentation['author_id'] != g.user['id']:
        return jsonify({'ok': False, 'error': 'Non autorizzato.'}), 403
    if slide_repository.move_slide_up(slide_id):
        slides = slide_repository.get_slides_by_presentation_id(slide['presentation_id'])
        return jsonify({'ok': True, 'slides': slides})
    return jsonify({'ok': False, 'error': 'Impossibile spostare la slide.'}), 400


@bp.route('/slides/<int:slide_id>/move_down', methods=['POST'])
@login_required
def move_slide_down(slide_id):
    slide = slide_repository.get_slide_by_id(slide_id)
    if not slide:
        return jsonify({'ok': False, 'error': 'Slide non trovata.'}), 404
    presentation = presentation_repository.get_presentation_by_id(slide['presentation_id'])
    if presentation['author_id'] != g.user['id']:
        return jsonify({'ok': False, 'error': 'Non autorizzato.'}), 403
    if slide_repository.move_slide_down(slide_id):
        slides = slide_repository.get_slides_by_presentation_id(slide['presentation_id'])
        return jsonify({'ok': True, 'slides': slides})
    return jsonify({'ok': False, 'error': 'Impossibile spostare la slide.'}), 400


@bp.route('/slides/<int:slide_id>', methods=['DELETE'])
@login_required
def delete_slide(slide_id):
    slide = slide_repository.get_slide_by_id(slide_id)
    if not slide:
        return jsonify({'ok': False, 'error': 'Slide non trovata.'}), 404
    presentation = presentation_repository.get_presentation_by_id(slide['presentation_id'])
    if presentation['author_id'] != g.user['id']:
        return jsonify({'ok': False, 'error': 'Non autorizzato.'}), 403
    slide_repository.delete_slide(slide_id)
    return jsonify({'ok': True, 'presentation_id': slide['presentation_id']})


# ─── Components ───────────────────────────────────────────────────────────────

@bp.route('/slides/<int:slide_id>/components/save', methods=['POST'])
@login_required
def save_slide_components(slide_id):
    slide = slide_repository.get_slide_by_id(slide_id)
    if not slide:
        return jsonify({'ok': False, 'error': 'Slide non trovata.'}), 404
    presentation = presentation_repository.get_presentation_by_id(slide['presentation_id'])
    if presentation['author_id'] != g.user['id']:
        return jsonify({'ok': False, 'error': 'Non autorizzato.'}), 403
    data = request.get_json() or {}
    bg_color = data.get('bg_color', slide.get('bg_color', '#ffffff'))
    slide_repository.update_slide(slide_id, slide['position'], bg_color)
    saved = slide_component_repository.save_all_components(slide_id, data.get('components', []))
    return jsonify({'ok': True, 'components': saved})


@bp.route('/slides/<int:slide_id>/upload_image', methods=['POST'])
@login_required
def upload_component_image(slide_id):
    slide = slide_repository.get_slide_by_id(slide_id)
    if not slide:
        return jsonify({'ok': False, 'error': 'Slide non trovata.'}), 404
    presentation = presentation_repository.get_presentation_by_id(slide['presentation_id'])
    if presentation['author_id'] != g.user['id']:
        return jsonify({'ok': False, 'error': 'Non autorizzato.'}), 403
    image_path = save_image(request.files.get('image'))
    if not image_path:
        return jsonify({'ok': False, 'error': 'Nessuna immagine caricata.'}), 400
    return jsonify({'ok': True, 'path': image_path})


# ─── Export / Import ──────────────────────────────────────────────────────────

@bp.route('/presentations/<int:presentation_id>/export.pptx')
@login_required
def export_pptx(presentation_id):
    from pptx import Presentation as PptxPres
    from pptx.util import Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN

    presentation = presentation_repository.get_presentation_by_id(presentation_id)
    if not presentation or presentation['author_id'] != g.user['id']:
        return jsonify({'ok': False, 'error': _('Presentazione non trovata.')}), 404

    W = 9144000
    H = 5143500
    prs = PptxPres()
    prs.slide_width = W
    prs.slide_height = H

    slides = slide_repository.get_slides_by_presentation_id(presentation_id)
    for slide_data in slides:
        pptx_slide = prs.slides.add_slide(prs.slide_layouts[6])

        bg = (slide_data.get('bg_color') or '#ffffff')
        if bg != 'transparent':
            fill = pptx_slide.background.fill
            fill.solid()
            hex_c = bg.lstrip('#')
            try:
                fill.fore_color.rgb = RGBColor(
                    int(hex_c[0:2], 16), int(hex_c[2:4], 16), int(hex_c[4:6], 16)
                )
            except Exception:
                pass

        components = slide_component_repository.get_components_by_slide(slide_data['id'])
        for comp in components:
            x = int(comp['x'] / 100 * W)
            y = int(comp['y'] / 100 * H)
            w = int(comp['width'] / 100 * W)
            h = int(comp['height'] / 100 * H)
            if comp['type'] in ('title', 'text', 'link'):
                txBox = pptx_slide.shapes.add_textbox(x, y, w, h)
                tf = txBox.text_frame
                tf.word_wrap = True
                tf.auto_size = MSO_AUTO_SIZE.NONE
                tf.margin_left = 0
                tf.margin_right = 0
                tf.margin_top = 0
                tf.margin_bottom = 0

                color = (comp.get('color') or '#333333')
                hex_c = color.lstrip('#') if color.startswith('#') else '333333'
                try:
                    rgb = RGBColor(int(hex_c[0:2], 16), int(hex_c[2:4], 16), int(hex_c[4:6], 16))
                except Exception:
                    rgb = RGBColor(0x33, 0x33, 0x33)

                font_size = comp.get('font_size') or 24
                pt_size = max(8, int(font_size * 0.72))

                raw_text = comp.get('content', '') or ''
                lines = raw_text.split('\n')

                first_para = tf.paragraphs[0]
                first_para.alignment = PP_ALIGN.LEFT
                run = first_para.add_run()
                run.text = lines[0]
                run.font.size = Pt(pt_size)
                run.font.bold = (comp['type'] == 'title')
                run.font.color.rgb = rgb

                for line in lines[1:]:
                    p = tf.add_paragraph()
                    p.alignment = PP_ALIGN.LEFT
                    r = p.add_run()
                    r.text = line
                    r.font.size = Pt(pt_size)
                    r.font.bold = (comp['type'] == 'title')
                    r.font.color.rgb = rgb
            elif comp['type'] == 'image':
                if comp.get('image'):
                    img_path = os.path.join(current_app.root_path, comp['image'].lstrip('/'))
                    if os.path.exists(img_path):
                        try:
                            pptx_slide.shapes.add_picture(img_path, x, y, w, h)
                        except Exception:
                            pass
                else:
                    from pptx.enum.shapes import MSO_SHAPE
                    box = pptx_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
                    box.fill.solid()
                    box.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
                    box.line.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
                    box.line.width = Pt(1)
                    box.shadow.inherit = False

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    title = presentation.get('title', 'presentation')
    return send_file(
        buf,
        download_name=f'{title}.pptx',
        as_attachment=True,
        mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
    )


@bp.route('/presentations/<int:presentation_id>/import', methods=['POST'])
@login_required
def import_pptx(presentation_id):
    from pptx import Presentation as PptxPres
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    presentation = presentation_repository.get_presentation_by_id(presentation_id)
    if not presentation or presentation['author_id'] != g.user['id']:
        return jsonify({'ok': False, 'error': _('Presentazione non trovata.')}), 404

    file = request.files.get('file')
    if not file or not file.filename.lower().endswith('.pptx'):
        return jsonify({'ok': False, 'error': 'File .pptx richiesto.'}), 400

    prs = PptxPres(file)
    W = int(prs.slide_width)
    H = int(prs.slide_height)
    if not W or not H:
        return jsonify({'ok': False, 'error': 'Dimensioni slide non valide.'}), 400

    existing = slide_repository.get_slides_by_presentation_id(presentation_id)
    position = len(existing) + 1
    count = 0

    for pptx_slide in prs.slides:
        bg_color = '#ffffff'
        try:
            fill = pptx_slide.background.fill
            if fill.type is not None:
                rgb = fill.fore_color.rgb
                bg_color = '#{:02x}{:02x}{:02x}'.format(rgb.red, rgb.green, rgb.blue)
        except Exception:
            pass

        slide = slide_repository.create_slide(presentation_id, position, bg_color)
        position += 1
        count += 1
        components = []

        for idx, shape in enumerate(pptx_slide.shapes):
            if shape.has_text_frame and shape.text_frame.text.strip():
                comp_type = 'text'
                try:
                    if shape.placeholder_format is not None and (
                        shape.placeholder_format.idx == 0 or 'title' in shape.name.lower()
                    ):
                        comp_type = 'title'
                except Exception:
                    pass
                if 'title' in shape.name.lower():
                    comp_type = 'title'

                font_size = 24
                color = '#333333'
                try:
                    for para in shape.text_frame.paragraphs:
                        if para.runs:
                            run = para.runs[0]
                            if run.font.size:
                                font_size = int(run.font.size.pt)
                            if run.font.color.type is not None:
                                rgb = run.font.color.rgb
                                color = '#{:02x}{:02x}{:02x}'.format(
                                    rgb.red, rgb.green, rgb.blue
                                )
                            break
                except Exception:
                    pass

                components.append({
                    'type': comp_type,
                    'content': shape.text_frame.text,
                    'x': round(int(shape.left) / W * 100, 1),
                    'y': round(int(shape.top) / H * 100, 1),
                    'width': round(int(shape.width) / W * 100, 1),
                    'height': round(int(shape.height) / H * 100, 1),
                    'font_size': font_size,
                    'color': color,
                    'bg_color': 'transparent',
                    'z_index': idx,
                    'image': None,
                })
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img_bytes = shape.image.blob
                    ext = shape.image.ext
                    fname = f"{uuid.uuid4().hex}.{ext}"
                    upload_folder = current_app.config['UPLOAD_FOLDER']
                    os.makedirs(upload_folder, exist_ok=True)
                    with open(os.path.join(upload_folder, fname), 'wb') as img_f:
                        img_f.write(img_bytes)
                    components.append({
                        'type': 'image',
                        'content': '',
                        'x': round(int(shape.left) / W * 100, 1),
                        'y': round(int(shape.top) / H * 100, 1),
                        'width': round(int(shape.width) / W * 100, 1),
                        'height': round(int(shape.height) / H * 100, 1),
                        'font_size': 0,
                        'color': '#000000',
                        'bg_color': 'transparent',
                        'z_index': idx,
                        'image': f'/static/uploads/{fname}',
                    })
                except Exception:
                    pass

        slide_component_repository.save_all_components(slide['id'], components)

    return jsonify({'ok': True, 'slides_added': count})


# ─── Auth ─────────────────────────────────────────────────────────────────────

@bp.route('/auth/login', methods=['POST'])
def login():
    data = request.get_json()
    username = data.get('username', '').strip()
    password = data.get('password', '').strip()
    from app.repositories import user_repository
    user = user_repository.authenticate_user(username, password)
    if user:
        from flask import session
        session.clear()
        session['user_id'] = user['id']
        return jsonify({'ok': True, 'username': username})
    return jsonify({'ok': False, 'error': 'Credenziali errate.'}), 401


@bp.route('/auth/register', methods=['POST'])
def register():
    data = request.get_json()
    username = data.get('username', '').strip()
    password = data.get('password', '').strip()
    email = data.get('email', '').strip()
    if not username or not password or not email:
        return jsonify({'ok': False, 'error': _('Username, email e password obbligatori.')}), 400
    if '@' not in email or '.' not in email:
        return jsonify({'ok': False, 'error': _('Inserisci un indirizzo email valido.')}), 400
    from werkzeug.security import generate_password_hash
    from app.repositories import user_repository
    try:
        hashed_pwd = generate_password_hash(password)
        user_repository.create_user(username, hashed_pwd, email)
        return jsonify({'ok': True})
    except ValueError as e:
        return jsonify({'ok': False, 'error': str(e)}), 400


@bp.route('/auth/logout', methods=['POST'])
def logout():
    from flask import session
    session.clear()
    return jsonify({'ok': True})
